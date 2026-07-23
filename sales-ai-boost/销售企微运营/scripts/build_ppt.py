# -*- coding: utf-8 -*-
"""生成《企微每周话题话术引擎·周会汇报》PPT — 通用版（无品牌、可套任意行业）。
原则：铺满无空白 / 板块内文字充实 / 正文>=15pt / 5页 / 不孤字 / 多色协调。"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- 主题色（通用红作强调色，不绑定任何品牌）----
BRAND_RED = RGBColor(0xE6, 0x00, 0x12)
DARK = RGBColor(0x26, 0x2A, 0x33)
GRAY = RGBColor(0x6B, 0x70, 0x78)
LIGHT_GRAY = RGBColor(0xF3, 0xF4, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED_BG = RGBColor(0xFD, 0xEC, 0xED)
BLUE_BG = RGBColor(0xEA, 0xF0, 0xF8)
CARD_BORDER = RGBColor(0xE4, 0xE6, 0xEB)
FONT = "微软雅黑"

# ---- 协调多色系（8类话题各一色，全篇复用）----
C_RED = RGBColor(0xE6, 0x00, 0x12)      # 1 公司好消息
C_BLUE = RGBColor(0x2B, 0x57, 0xA0)     # 2 行业趋势/政策
C_TEAL = RGBColor(0x0E, 0x8A, 0x7D)     # 3 客户专属喜讯
C_AMBER = RGBColor(0xE0, 0x7B, 0x12)    # 4 政策补贴/以旧换新
C_PURPLE = RGBColor(0x6B, 0x4F, 0xA0)   # 5 实用干货
C_GREEN = RGBColor(0x1E, 0x7A, 0x3C)    # 6 节日/节气/关怀
C_MAGENTA = RGBColor(0xC2, 0x18, 0x6A)  # 7 公司活动/权益
C_INDIGO = RGBColor(0x3F, 0x51, 0xB5)   # 8 AI/算力热点
CAT_COLORS = [C_RED, C_BLUE, C_TEAL, C_AMBER, C_PURPLE, C_GREEN, C_MAGENTA, C_INDIGO]
SEG_COLORS = [C_RED, C_BLUE, C_TEAL, C_AMBER, C_GREEN]

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]
TOTAL = 5


def add_slide():
    return prs.slides.add_slide(BLANK)


def rect(slide, x, y, w, h, fill=None, line=None, line_w=0.75, round_=False):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def _set_font(run, size, color=DARK, bold=False, font=FONT):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font
    rPr = run._r.get_or_add_rPr()
    rPr.set('{http://schemas.openxmlformats.org/drawingml/2006/main}altLang', 'en-US')
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', font)


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(3); tf.margin_right = Pt(3)
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    return tb, tf


def para(tf, text, size=15, color=DARK, bold=False, align=PP_ALIGN.LEFT,
         space_after=6, space_before=0, first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    p.level = 0
    r = p.add_run(); r.text = text
    _set_font(r, size, color, bold)
    return p


def header_band(slide, title, kicker=None):
    rect(slide, 0, 0, SW, Inches(1.1), fill=BRAND_RED)
    rect(slide, 0, Inches(1.1), SW, Pt(3), fill=DARK)
    tb, tf = textbox(slide, Inches(0.5), Inches(0.12), Inches(12.3), Inches(0.9),
                     anchor=MSO_ANCHOR.MIDDLE)
    if kicker:
        para(tf, kicker, size=12, color=WHITE, bold=True, first=True, space_after=2)
        para(tf, title, size=25, color=WHITE, bold=True, space_after=0)
    else:
        para(tf, title, size=28, color=WHITE, bold=True, first=True)


def footer(slide, page):
    tb, tf = textbox(slide, Inches(0.5), Inches(7.08), Inches(10), Inches(0.34))
    para(tf, "企微每周话题话术引擎 · 部门周会汇报", size=11, color=GRAY, first=True)
    tb2, tf2 = textbox(slide, Inches(11.8), Inches(7.08), Inches(1.2), Inches(0.34))
    para(tf2, f"{page}/{TOTAL}", size=11, color=GRAY, first=True, align=PP_ALIGN.RIGHT)


def table(slide, x, y, w, h, rows, col_w, header_fill=BRAND_RED,
          font_size=14.5, header_size=16, row_h=Inches(0.56),
          col0_colors=None):
    n_rows = len(rows); n_cols = len(rows[0])
    gtbl = slide.shapes.add_table(n_rows, n_cols, x, y, w, h).table
    gtbl.first_row = False
    gtbl.horz_banding = False
    total = sum(col_w)
    for j, cw in enumerate(col_w):
        gtbl.columns[j].width = Emu(int(w * cw / total))
    for i in range(n_rows):
        gtbl.rows[i].height = row_h
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = gtbl.cell(i, j)
            cell.margin_left = Pt(6); cell.margin_right = Pt(5)
            cell.margin_top = Pt(2); cell.margin_bottom = Pt(2)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if i == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = header_fill
                tcol, tbold, tsz = WHITE, True, header_size
            elif j == 0 and col0_colors is not None:
                cell.fill.solid(); cell.fill.fore_color.rgb = col0_colors[i - 1]
                tcol, tbold, tsz = WHITE, True, font_size
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if i % 2 else LIGHT_GRAY
                tcol, tbold, tsz = DARK, False, font_size
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run(); r.text = str(val)
            _set_font(r, tsz, tcol, tbold)
    return gtbl


# ============================================================
# 第1页：封面
# ============================================================
s = add_slide()
rect(s, 0, 0, SW, SH, fill=WHITE)
rect(s, 0, 0, Inches(0.5), SH, fill=BRAND_RED)
tb, tf = textbox(s, Inches(0.9), Inches(1.3), Inches(7.6), Inches(2.6))
para(tf, "用 AI 助力企微运营", size=46, color=DARK, bold=True, first=True, space_after=4)
para(tf, "每周话题话术引擎", size=46, color=BRAND_RED, bold=True, space_after=0)
tb, tf = textbox(s, Inches(0.95), Inches(4.0), Inches(7.7), Inches(2.0))
para(tf, "每周一 5 分钟，拿到 5 条「可发、不尬、能挖商机」的企微话术",
     size=20, color=GRAY, first=True, space_after=12)
para(tf, "1 个引擎自动跑　·　销售只复制粘贴　·　内容统一　·　红线自检",
     size=16, color=DARK, bold=True, space_after=0)
rect(s, Inches(8.9), Inches(1.3), Inches(3.95), Inches(4.7), fill=BRAND_RED, round_=True)
tb, tf = textbox(s, Inches(9.1), Inches(1.7), Inches(3.6), Inches(4.0),
                 anchor=MSO_ANCHOR.MIDDLE)
para(tf, "中心化生产", size=30, color=WHITE, bold=True, first=True, align=PP_ALIGN.CENTER, space_after=6)
para(tf, "+", size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER, space_after=6)
para(tf, "全员分发", size=30, color=WHITE, bold=True, align=PP_ALIGN.CENTER, space_after=14)
para(tf, "零部署", size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER, space_after=4)
para(tf, "内容统一", size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER, space_after=4)
para(tf, "每周自动产出", size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER, space_after=0)
rect(s, Inches(0.9), Inches(6.45), Inches(11.9), Inches(0.62), fill=LIGHT_GRAY, round_=True)
tb, tf = textbox(s, Inches(1.1), Inches(6.45), Inches(11.5), Inches(0.62), anchor=MSO_ANCHOR.MIDDLE)
para(tf, "B2B 销售团队 · 企微运营提效 · 部门周会汇报 · 2026-07-22",
     size=14, color=DARK, bold=True, first=True)

# ============================================================
# 第2页：痛点 → 解法
# ============================================================
s = add_slide()
header_band(s, "为什么做 · 怎么落地", kicker="一图看清：痛点与解法")
panel_y, panel_h = Inches(1.35), Inches(5.55)
header_h = Inches(0.6)


def panel(side, header_text, hcolor, hbg, items):
    px = Inches(0.5) if side == "L" else Inches(6.98)
    pw = Inches(5.85)
    rect(s, px, panel_y, pw, panel_h, fill=hbg, round_=True)
    rect(s, px, panel_y, pw, header_h, fill=hcolor, round_=True)
    rect(s, px, panel_y + Inches(0.42), pw, Inches(0.18), fill=hcolor)
    tb, tf = textbox(s, px + Inches(0.25), panel_y, pw - Inches(0.4), header_h,
                     anchor=MSO_ANCHOR.MIDDLE)
    para(tf, header_text, size=20, color=WHITE, bold=True, first=True)
    n = len(items)
    top = panel_y + header_h + Inches(0.18)
    area = panel_h - header_h - Inches(0.3)
    ih = Emu(int(area / n))
    for i, (title, desc) in enumerate(items):
        y = top + ih * i
        rect(s, px + Inches(0.28), y + Emu(int(ih * 0.18)), Inches(0.16), Inches(0.16),
             fill=hcolor, round_=True)
        tb, tf = textbox(s, px + Inches(0.6), y, pw - Inches(0.9), ih,
                         anchor=MSO_ANCHOR.MIDDLE)
        para(tf, title, size=16.5, color=hcolor, bold=True, first=True, space_after=3)
        para(tf, desc, size=14, color=DARK, space_after=0)

panel("L", "现在的痛点", BRAND_RED, RED_BG, [
    ("不知道发什么", "群发通稿像骚扰，客户不回、关系变僵"),
    ("素材散落各处", "新闻、政策、补贴没人汇总成现成话术"),
    ("部署推广很难", "每人装 AI 工具，成本高、内容不统一"),
    ("关怀变成减分", "没话术、没承接，发了等于白发"),
])
panel("R", "我们的解法", C_BLUE, BLUE_BG, [
    ("一个引擎自动跑", "每周一 09:00 联网抓取最新素材"),
    ("八类话题矩阵", "筛 5 个话题，生成带完整话术卡"),
    ("销售零部署", "只复制粘贴发企微，内容统一"),
    ("四条红线自检", "不编造、不承诺价、单信息、可溯源"),
    ("已经上线运行", "自动化已建，本周实样已产出"),
])
circle = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.32), Inches(3.95), Inches(0.68), Inches(0.68))
circle.fill.solid(); circle.fill.fore_color.rgb = DARK; circle.line.fill.background()
ctb, ctf = textbox(s, Inches(6.32), Inches(3.95), Inches(0.68), Inches(0.68), anchor=MSO_ANCHOR.MIDDLE)
para(ctf, "→", size=26, color=WHITE, bold=True, first=True, align=PP_ALIGN.CENTER)
footer(s, 2)

# ============================================================
# 第3页：8 类话题矩阵
# ============================================================
s = add_slide()
header_band(s, "核心方法论 ①：话题 8 类矩阵", kicker="从「关怀入口」到「商机承接」的完整链路")
rows = [
    ["#", "话题类别", "发出目的", "公司衔接点", "承接的商机"],
    ["1", "公司好消息", "证明实力、引入合作", "新品 / 技术 / 全栈能力", "产品切入"],
    ["2", "行业趋势 / 政策", "问需求、显专业", "数字化 / AI 方案", "数字化 · 算力"],
    ["3", "客户企业专属喜讯", "恭喜、拉近关系（须真实）", "扩产要设备跟上", "扩产采购"],
    ["4", "政策补贴 / 以旧换新", "「有笔钱你能省」", "批量换新 / 订阅式服务", "批量换新"],
    ["5", "实用干货 tips", "立专家人设", "可管理 / 安全 / 降本", "长期信任"],
    ["6", "节日 / 节气 / 关怀", "纯暖场、无压力", "不推产品", "维系关系"],
    ["7", "公司活动 / 权益", "直接转化钩子", "活动 / 优惠 / 会员", "报名 · 询价"],
    ["8", "AI / 算力热点借势", "蹭趋势、聊痛点", "AI PC / 工作站 / 算力", "工作站 · 算力"],
]
table(s, Inches(0.5), Inches(1.35), Inches(12.33), Inches(5.55), rows,
      col_w=[0.7, 2.9, 4.2, 2.9, 3.0], font_size=15, header_size=16,
      row_h=Inches(0.56), col0_colors=CAT_COLORS)
footer(s, 3)

# ============================================================
# 第4页：五段式话术公式 + 四条红线 + 范例
# ============================================================
s = add_slide()
header_band(s, "核心方法论 ②：话术怎么写", kicker="五段式公式 + 四条红线 + 实样")
rect(s, Inches(0.5), Inches(1.3), Inches(12.33), Inches(0.78), fill=LIGHT_GRAY, round_=True)
seg_names = ["开场关怀", "事件钩子", "公司衔接", "价值锚点", "轻 CTA"]
seg_w = Inches(2.18); gap = Inches(0.13); x0 = Inches(0.62); y = Inches(1.42)
for i, nm in enumerate(seg_names):
    x = x0 + (seg_w + gap) * i
    rect(s, x, y, seg_w, Inches(0.54), fill=SEG_COLORS[i], round_=True)
    tb, tf = textbox(s, x, y, seg_w, Inches(0.54), anchor=MSO_ANCHOR.MIDDLE)
    para(tf, nm, size=15, color=WHITE, bold=True, first=True, align=PP_ALIGN.CENTER)
tb, tf = textbox(s, Inches(0.5), Inches(2.35), Inches(6.0), Inches(0.45))
para(tf, "五段拆解", size=18, color=DARK, bold=True, first=True)
tb, tf = textbox(s, Inches(0.5), Inches(2.85), Inches(6.05), Inches(2.05))
seg_desc = [
    ("开场关怀", "称呼 + 一句人情味问候"),
    ("事件钩子", "本周话题，带时效感（「刚」「最新」）"),
    ("公司衔接", "自然过渡到能力，不硬广"),
    ("价值锚点", "只推一个（TCO更省 / 数据不出厂）"),
    ("轻 CTA", "低压力下一步（「发份资料参考?」）"),
]
for i, (t, d) in enumerate(seg_desc):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(7)
    rk = p.add_run(); rk.text = "● "; _set_font(rk, 15, SEG_COLORS[i], True)
    rt = p.add_run(); rt.text = t + "："; _set_font(rt, 15, DARK, True)
    rd = p.add_run(); rd.text = d; _set_font(rd, 15, DARK)
tb, tf = textbox(s, Inches(0.5), Inches(5.0), Inches(6.0), Inches(0.45))
para(tf, "四条红线（每次必过）", size=18, color=BRAND_RED, bold=True, first=True)
tb, tf = textbox(s, Inches(0.5), Inches(5.5), Inches(6.05), Inches(1.5))
red_lines = [
    "不编造客户事件：专属喜讯须真实可核实来源",
    "不承诺具体价格：只说「有批量方案 / 专属价」",
    "单信息原则：一条话术只推一个核心信息",
    "来源可溯：素材留链接；关怀类只暖场不硬推",
]
for i, t in enumerate(red_lines):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(6)
    rk = p.add_run(); rk.text = "● "; _set_font(rk, 15, BRAND_RED, True)
    rd = p.add_run(); rd.text = t; _set_font(rd, 15, DARK)
rect(s, Inches(6.75), Inches(2.35), Inches(6.08), Inches(4.65), fill=LIGHT_GRAY, round_=True)
rect(s, Inches(6.75), Inches(2.35), Inches(6.08), Inches(0.55), fill=C_RED, round_=True)
rect(s, Inches(6.75), Inches(2.72), Inches(6.08), Inches(0.18), fill=C_RED)
tb, tf = textbox(s, Inches(6.95), Inches(2.35), Inches(5.7), Inches(0.55), anchor=MSO_ANCHOR.MIDDLE)
para(tf, "实样话术（类别①公司好消息）", size=16, color=WHITE, bold=True, first=True)
tb, tf = textbox(s, Inches(6.95), Inches(3.05), Inches(5.7), Inches(3.0))
example = ("X总早上好，分享个好消息：我们上半年业绩创了新高，还刚入选了年度行业解决方案榜单。"
           "其实我们早就不只是卖单品，从员工的终端到机房服务器、再到产线 AI 方案都能接。"
           "您这边要是有 IT 或研发设备的规划，随时喊我帮您把方案捋一捋。")
para(tf, example, size=15.5, color=DARK, first=True, space_after=10)
para(tf, "承接动作：客户表示「在看 / 在规划」→ 约需求盘点，问清岗位、软件、负载。",
     size=14, color=GRAY, space_after=0)
footer(s, 4)

# ============================================================
# 第5页：实样 Demo + 落地 SOP
# ============================================================
s = add_slide()
header_band(s, "实样 Demo + 怎么跑起来", kicker="本周话题包示例 · 直接复制发企微")
demo = [
    ("① 政策补贴 / 以旧换新", C_AMBER,
     "2026「两新」政策落地，数码及智能产品购新补贴",
     "X总好，今年国家「两新」政策力度不小——大规模设备更新 + 数码和智能产品购新补贴都在推，中小企业申报门槛还降了。您那边要是有用了好几年的老旧电脑或工作站，正好顺着这波批量换掉。我帮您粗算旧机折抵加政策能省多少，您看要不要？",
     "承接：客户有意向 → 盘点在用老旧设备，做换新 + 订阅式服务两套方案"),
    ("② AI / 算力热点", C_INDIGO,
     "新款 AI 工作站，本地跑 32B 大模型",
     "X总好，现在 AI 这么火，制造业先落地的往往是研发端。我们刚上市新款 AI 工作站，本地就能稳跑 32B 大模型，做仿真、建模、AI 开发都不用挤云上、数据也不出厂。您们研发跑大型图纸或仿真现在扛得住吗？",
     "承接：客户吐槽卡 / 崩 → 索取软件清单与负载，出配置建议"),
    ("③ 节气关怀", C_GREEN,
     "大暑时节，高温关怀（纯暖场，不推产品）",
     "X总，大暑了，这两天是真晒。车间和现场注意通风降温，提醒大家多喝水别中暑，周末也记得歇歇。有需要随时招呼我，不打扰您了。",
     "承接：本条不推销，回暖回应即达目的"),
]
x0 = Inches(0.5); col_w = Inches(4.0); gap = Inches(0.165); y0 = Inches(1.32)
for i, (tag, color, title, speech, action) in enumerate(demo):
    x = x0 + (col_w + gap) * i
    rect(s, x, y0, col_w, Inches(4.15), fill=WHITE, line=CARD_BORDER, line_w=1.25, round_=True)
    rect(s, x, y0, col_w, Inches(0.55), fill=color, round_=True)
    rect(s, x, y0 + Inches(0.37), col_w, Inches(0.18), fill=color)
    tb, tf = textbox(s, x + Inches(0.1), y0, col_w - Inches(0.2), Inches(0.55), anchor=MSO_ANCHOR.MIDDLE)
    para(tf, tag, size=15, color=WHITE, bold=True, first=True, align=PP_ALIGN.CENTER)
    tb, tf = textbox(s, x + Inches(0.2), y0 + Inches(0.68), col_w - Inches(0.4), Inches(1.0))
    para(tf, title, size=14.5, color=DARK, bold=True, first=True, space_after=0)
    rect(s, x + Inches(0.2), y0 + Inches(1.7), col_w - Inches(0.4), Inches(1.95), fill=LIGHT_GRAY, round_=True)
    tb, tf = textbox(s, x + Inches(0.3), y0 + Inches(1.78), col_w - Inches(0.6), Inches(1.8), anchor=MSO_ANCHOR.TOP)
    para(tf, speech, size=13.5, color=DARK, first=True, space_after=0)
    tb, tf = textbox(s, x + Inches(0.2), y0 + Inches(3.72), col_w - Inches(0.4), Inches(0.4))
    para(tf, action, size=13, color=GRAY, first=True, space_after=0)
rect(s, Inches(0.5), Inches(5.7), Inches(12.33), Inches(1.2), fill=BRAND_RED, round_=True)
tb, tf = textbox(s, Inches(0.7), Inches(5.78), Inches(12.0), Inches(0.45))
para(tf, "落地 SOP（每周一 09:00 自动产出，销售只复制粘贴）", size=15, color=WHITE, bold=True, first=True)
tb, tf = textbox(s, Inches(0.7), Inches(6.22), Inches(12.0), Inches(0.65))
para(tf, "① 引擎自动产 5 条话题卡　→　② 挑合适话题发对应客户　→　③ 换真实称呼发企微　→　④ 按承接动作接住别断在暖场　→　⑤ 月底统计回复率迭代话术库",
     size=13.5, color=WHITE, first=True, space_after=0)
footer(s, 5)

out = "示例/企微话题引擎_周会汇报.pptx"
prs.save(out)
print("saved:", out, "slides:", len(prs.slides._sldIdLst))
